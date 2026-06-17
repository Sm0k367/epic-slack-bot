#!/usr/bin/env python3
"""
Epic Agent Ultimate God Mode v6
Complete working Slack bot - all skills included.
Set tokens via environment variables before running.
"""
import os
import sys
import time
import logging
import requests
import tempfile
import subprocess
from collections import defaultdict, deque
from slack_bolt import App
from slack_bolt.adapter.socket_mode import SocketModeHandler

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("EpicAgent")

# === LOAD TOKENS FROM ENVIRONMENT ===
SLACK_BOT_TOKEN = os.environ.get("SLACK_BOT_TOKEN")
SLACK_APP_TOKEN = os.environ.get("SLACK_APP_TOKEN")
XAI_KEY = os.environ.get("XAI_API_KEY", "")
PIXIO_KEY = os.environ.get("PIXIO_API_KEY", "")

if not SLACK_BOT_TOKEN or not SLACK_APP_TOKEN:
    print("""
╔════════════════════════════════════════════════════════════╗
║  Epic Agent Ultimate v6 - Missing Tokens                   ║
╚════════════════════════════════════════════════════════════╝

Please set your Slack tokens first:

export SLACK_BOT_TOKEN="xoxb-11370202537843-..."
export SLACK_APP_TOKEN="xapp-1-A0BAZL0LCBX-..."

Then run:
python3 EpicAgent_ULTIMATE.py
""")
    sys.exit(1)

app = App(token=SLACK_BOT_TOKEN)
HISTORY = defaultdict(lambda: deque(maxlen=15))

def add_mem(ch, user, text):
    HISTORY[ch].append({"user": user, "text": text})

def get_mem(ch):
    return list(HISTORY[ch])

def send_long(say, text, thread_ts):
    if len(text) <= 3500:
        say(text=text, thread_ts=thread_ts)
        return
    chunks, current = [], ""
    for line in text.split("\n"):
        if len(current) + len(line) + 1 > 3500:
            if current: chunks.append(current)
            current = line + "\n"
        else:
            current += line + "\n"
    if current: chunks.append(current)
    for i, c in enumerate(chunks):
        say(text=(f"[{i+1}/{len(chunks)}] " if len(chunks) > 1 else "") + c, thread_ts=thread_ts)
        time.sleep(0.25)

def grok(text, user, ch, extra=""):
    ctx = get_mem(ch)
    msgs = [{"role": "system", "content": f"You are Epic Agent Ultimate v6. {extra}"}]
    for m in ctx[-12:]:
        msgs.append({"role": "user", "content": f"{m['user']}: {m['text']}"})
    msgs.append({"role": "user", "content": text})
    try:
        r = requests.post("https://api.x.ai/v1/chat/completions",
            headers={"Authorization": f"Bearer {XAI_KEY}", "Content-Type": "application/json"},
            json={"model": "grok-3", "messages": msgs, "max_tokens": 4000}, timeout=30)
        return r.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"Error: {e}"

def pixio_image(prompt):
    if not PIXIO_KEY:
        return None
    try:
        models = requests.get("https://beta.pixio.myapps.ai/api/v1/models",
            headers={"Authorization": f"Bearer {PIXIO_KEY}"}, timeout=8).json()
        model = next((m["id"] for m in models.get("models", []) if "flux" in m.get("id", "").lower()), "pixio/flux/dev")
        gen = requests.post("https://beta.pixio.myapps.ai/api/v1/generate",
            headers={"Authorization": f"Bearer {PIXIO_KEY}", "Content-Type": "application/json"},
            json={"providerId": "pixio", "modelId": model, "params": {"prompt": prompt}}, timeout=12).json()
        cid = gen.get("contentId") or gen.get("id")
        for _ in range(40):
            time.sleep(2)
            st = requests.get(f"https://beta.pixio.myapps.ai/api/v1/generations/{cid}",
                headers={"Authorization": f"Bearer {PIXIO_KEY}"}, timeout=6).json()
            if st.get("status") == "succeeded":
                return st.get("outputUrl") or st.get("outputs", {}).get("imageUrl")
            if st.get("status") == "failed":
                return None
        return None
    except:
        return None

def skill_pdf(text, title="Report"):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        path = f"/tmp/epic_{int(time.time())}.pdf"
        c = canvas.Canvas(path, pagesize=letter)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, 750, title)
        c.setFont("Helvetica", 10)
        y = 720
        for line in text.split("\n"):
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(72, y, line[:95])
            y -= 13
        c.save()
        return path
    except Exception as e:
        return str(e)

def skill_pptx(title, bullets):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        path = f"/tmp/epic_{int(time.time())}.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        for b in bullets[:12]:
            p = tf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(16)
        prs.save(path)
        return path
    except Exception as e:
        return str(e)

def skill_docx(text, title):
    try:
        from docx import Document
        path = f"/tmp/epic_{int(time.time())}.docx"
        d = Document()
        d.add_heading(title, 0)
        d.add_paragraph(text)
        d.save(path)
        return path
    except Exception as e:
        return str(e)

def skill_code(code):
    try:
        with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
            f.write(code)
            path = f.name
        out = subprocess.run(["python3", path], capture_output=True, text=True, timeout=20)
        os.unlink(path)
        return out.stdout + out.stderr
    except Exception as e:
        return str(e)

def route(text, user, ch, say, client, thread_ts):
    add_mem(ch, user, text)
    t = text.lower()

    # Image generation
    if any(w in t for w in ["generate", "make", "create"]) and any(w in t for w in ["image", "pic", "photo", "logo"]):
        say(text="🎨 Generating image...", thread_ts=thread_ts)
        url = pixio_image(text)
        send_long(say, url or "Image generation failed", thread_ts)
        return

    # PDF
    if "pdf" in t or "report" in t:
        say(text="📄 Creating PDF report...", thread_ts=thread_ts)
        content = grok(f"Write a detailed report about: {text}", user, ch)
        path = skill_pdf(content, "Epic Report")
        say(text=f"✅ PDF created: {path}", thread_ts=thread_ts)
        return

    # PPTX
    if "pptx" in t or "presentation" in t or "slides" in t:
        say(text="📊 Creating presentation...", thread_ts=thread_ts)
        content = grok(f"Create presentation content about: {text}", user, ch)
        path = skill_pptx("Epic Presentation", content.split("\n")[:10])
        say(text=f"✅ PPTX created: {path}", thread_ts=thread_ts)
        return

    # DOCX
    if "docx" in t or "document" in t:
        say(text="📝 Creating document...", thread_ts=thread_ts)
        content = grok(f"Write a professional document about: {text}", user, ch)
        path = skill_docx(content, "Epic Document")
        say(text=f"✅ DOCX created: {path}", thread_ts=thread_ts)
        return

    # Code execution
    if "run code" in t or "execute" in t:
        code = text.split("```")[-2] if "```" in text else text
        say(text="💻 Executing code...", thread_ts=thread_ts)
        result = skill_code(code)
        send_long(say, f"```\n{result}\n```", thread_ts)
        return

    # Default chat with memory
    send_long(say, grok(text, user, ch), thread_ts)

@app.event("message")
def on_message(event, say, client):
    if event.get("bot_id"):
        return
    uid = event.get("user")
    name = client.users_info(user=uid)["user"].get("real_name", "there")
    txt = event.get("text", "")
    ch = event.get("channel")
    if ch == "C0BBWQVUAQ0" and len(txt) > 1:
        route(txt, name, ch, say, client, event.get("thread_ts") or event.get("ts"))
    elif event.get("channel_type") == "im":
        route(txt, name, ch, say, client, None)

if __name__ == "__main__":
    log.info("🚀 Epic Agent Ultimate God Mode v6 starting...")
    SocketModeHandler(app, SLACK_APP_TOKEN).start()
